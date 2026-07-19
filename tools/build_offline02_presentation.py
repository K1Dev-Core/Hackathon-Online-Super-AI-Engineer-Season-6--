from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = Path("/Users/k1god/Downloads/Black and Gray Gradient Professional Presentation.pptx")
OUTPUT = ROOT / "presentation" / "SuperAI6_IoT_Attack_Detection_Offline02_Presentation_TH.pptx"

WHITE = RGBColor(247, 247, 247)
MUTED = RGBColor(205, 205, 205)
FONT = "Noto Sans Thai"


def style_text(shape, text, size, *, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
               valign=MSO_ANCHOR.TOP, margins=(0.16, 0.10, 0.16, 0.10), line_spacing=1.05):
    """Replace a template text box while preserving its position and shape."""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margins[0])
    tf.margin_top = Inches(margins[1])
    tf.margin_right = Inches(margins[2])
    tf.margin_bottom = Inches(margins[3])
    tf.vertical_anchor = valign
    tf.text = text

    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        paragraph.line_spacing = line_spacing
        for run in paragraph.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def apply_slide(slide, *, title=None, body=None, small=None, narration=None,
                left=None, badge=None, tagline=None, title_size=31):
    if title is not None:
        style_text(slide.shapes[2], title, title_size, bold=True, valign=MSO_ANCHOR.MIDDLE)
    if body is not None:
        style_text(slide.shapes[3], body, 15.5, bold=True, color=WHITE, line_spacing=1.08)
    if small is not None:
        style_text(slide.shapes[4], small, 15, bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE)
    if narration is not None:
        # The narration box is the last large text area in the template. The index
        # differs on slides with an extra badge, so callers pass the right shape.
        pass
    if left is not None:
        style_text(slide.shapes[3], left, 15.5, bold=True, color=WHITE, line_spacing=1.08)
    if badge is not None:
        style_text(slide.shapes[5], badge, 12.5, bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE)
    if tagline is not None:
        style_text(slide.shapes[5], tagline, 11.5, bold=True, color=MUTED, valign=MSO_ANCHOR.MIDDLE)


def narration(slide, shape_index, text):
    style_text(
        slide.shapes[shape_index],
        text,
        13.2,
        color=MUTED,
        bold=False,
        margins=(0.16, 0.06, 0.16, 0.04),
        line_spacing=1.0,
    )


def add_panel(slide, left, top, width, height, title, sections, footer=None):
    """Cover a template illustration with a clean project-specific information card."""
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(18, 21, 26)
    panel.line.color.rgb = RGBColor(105, 110, 118)
    panel.line.width = Pt(1.2)

    title_box = slide.shapes.add_textbox(Inches(left + 0.42), Inches(top + 0.32), Inches(width - 0.84), Inches(0.55))
    style_text(title_box, title, 18, bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE)

    y = top + 1.08
    section_height = (height - 1.55) / max(1, len(sections))
    for heading, detail in sections:
        box = slide.shapes.add_textbox(Inches(left + 0.42), Inches(y), Inches(width - 0.84), Inches(section_height - 0.10))
        style_text(box, f"{heading}\n{detail}", 13.5, bold=False, color=WHITE, line_spacing=1.02)
        y += section_height

    if footer:
        footer_box = slide.shapes.add_textbox(Inches(left + 0.42), Inches(top + height - 0.50), Inches(width - 0.84), Inches(0.30))
        style_text(footer_box, footer, 10.5, bold=True, color=MUTED, valign=MSO_ANCHOR.MIDDLE)


def build():
    if not TEMPLATE.exists():
        raise FileNotFoundError(TEMPLATE)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(TEMPLATE))

    # Slide 1: opening
    slide = prs.slides[0]
    style_text(slide.shapes[2], "IOT ATTACK DETECTION\nTHROUGH TCP + MQTT", 31, bold=True, valign=MSO_ANCHOR.MIDDLE)
    style_text(slide.shapes[3], "แชร์วิธีคิดจากข้อมูลจริง", 16, bold=True, valign=MSO_ANCHOR.MIDDLE)
    style_text(slide.shapes[4], "Super AI Engineer Season 6  •  Technical Story  •  610686-วชิรวิทย์", 11.5, color=MUTED)
    style_text(slide.shapes[5], "AUTHOR : WACHIRAWIT 610686", 10.5, bold=True, color=MUTED, valign=MSO_ANCHOR.MIDDLE)

    # Slide 2: problem
    slide = prs.slides[1]
    style_text(slide.shapes[1], "THE PROBLEM\nIN ONE SENTENCE", 30, bold=True, valign=MSO_ANCHOR.MIDDLE)
    style_text(slide.shapes[2], "เราต้องดู Network Traffic 10,000 แถว\nแล้วตอบเพียง 2 ค่า\n\n0 = Normal\n1 = Attack\n\nแต่ Train 100,000 แถว\nมี Normal เท่านั้น", 15.5, bold=True, line_spacing=1.02)
    style_text(slide.shapes[4], "โจทย์จริง", 15, bold=True, valign=MSO_ANCHOR.MIDDLE)
    narration(slide, 5, "สิ่งที่ทำให้โจทย์นี้ยากคือเราไม่เคยเห็นตัวอย่าง Attack ใน Train เลยครับ เราจึงไม่สามารถสอนโมเดลแบบปกติว่า Attack หน้าตาเป็นอย่างไรได้ ต้องเริ่มจากการทำความเข้าใจว่า Normal มีรูปแบบแบบไหนก่อน")

    # Slide 3: data
    slide = prs.slides[2]
    style_text(slide.shapes[2], "DATA\nCHALLENGE", 31, bold=True, valign=MSO_ANCHOR.MIDDLE)
    style_text(slide.shapes[3], "TRAIN 100,000 rows\nNORMAL ONLY\n\nTEST 10,000 rows\nNORMAL + ATTACK\n\n22 network features + Id\nTCP / MQTT / timing", 15.5, bold=True, line_spacing=1.04)
    style_text(slide.shapes[4], "ข้อมูลที่ให้มา", 15, bold=True, valign=MSO_ANCHOR.MIDDLE)
    narration(slide, 5, "Train มีแต่ Normal ทำให้โจทย์นี้เป็น One-Class Anomaly Detection มากกว่า Binary Classification ครับ จุดสำคัญคือข้อมูลหนึ่งแถวไม่ควรถูกตัดสินจากคอลัมน์เดียว แต่ต้องดูความสัมพันธ์ของ TCP, MQTT และบริบทของ connection")

    # Slide 4: main idea
    slide = prs.slides[3]
    style_text(slide.shapes[2], "MAIN IDEA\nONE-CLASS DETECTION", 29, bold=True, valign=MSO_ANCHOR.MIDDLE)
    style_text(slide.shapes[3], "เปรียบเหมือน รปภ. ที่รู้จักพฤติกรรมปกติของคนในอาคาร แม้ไม่เคยเห็นผู้บุกรุกมาก่อน ก็ยังสังเกตพฤติกรรมที่แตกต่างได้", 15.2, bold=True, line_spacing=1.08)
    style_text(slide.shapes[4], "1  เรียนรู้ Normal profile\n2  เปรียบเทียบกับ Test\n3  อ่านความหมาย TCP + MQTT\n4  แจ้งเฉพาะความผิดปกติที่มีหลักฐาน", 15, bold=True, line_spacing=1.08)
    narration(slide, 5, "ผมไม่ได้ให้โมเดลจำว่า Attack ต้องมีหน้าตาแบบเดียว แต่ให้มันรู้จักขอบเขตของ Normal ก่อน จากนั้นจึงค่อยถามว่าแถวนี้แตกต่างเพราะเป็น Attack จริง หรือแค่เป็น Normal รูปแบบใหม่")

    # Slide 5: features
    slide = prs.slides[4]
    style_text(slide.shapes[2], "FEATURES\nTHAT MATTER", 31, bold=True, valign=MSO_ANCHOR.MIDDLE)
    style_text(slide.shapes[3], "TCP\n• frame.len / tcp.len\n• tcp.window_size\n• SYN / ACK / RST\n• tcp.stream\n\nMQTT\n• mqtt.msgtype\n• mqtt.kalive\n• mqtt.len / topic_len", 14.5, bold=True, line_spacing=1.02)
    narration(slide, 4, "Feature ที่ผมให้ความสำคัญมีสองกลุ่มครับ กลุ่มแรกคือ TCP ซึ่งบอกขนาด packet และ state ของ connection กลุ่มที่สองคือ MQTT ซึ่งบอกว่า packet นี้เป็น CONNECT, PUBLISH หรือ SUBSCRIBE การอ่านสองกลุ่มร่วมกันช่วยลดการตีความผิดจากค่าที่ดูแปลกเพียงตัวเดียว")

    # Slide 6: process
    slide = prs.slides[5]
    style_text(slide.shapes[2], "FROM RULES\nTO A ROBUST MODEL", 29, bold=True, valign=MSO_ANCHOR.MIDDLE)
    style_text(slide.shapes[3], "1  Normal profile\n   set ของ window และ packet shape ที่พบจริง\n\n2  Protocol rules\n   SYN flood / CONNECT / SUBSCRIBE / PUBLISH\n\n3  Context gate\n   ดู stream และ family ไม่ดู packet เดี่ยว\n\n4  Hard negative\n   ตัด PINGRESP ที่ดูแปลกแต่เป็น Normal", 13.7, bold=True, line_spacing=1.02)
    # This template's small label starts slightly outside the canvas; keep it
    # empty rather than letting the first characters be clipped.
    style_text(slide.shapes[4], "", 15, bold=True, valign=MSO_ANCHOR.MIDDLE)
    narration(slide, 5, "ผมเริ่มจากกฎที่อธิบายได้ก่อนครับ เพราะเรายังไม่มี Attack label ที่เชื่อถือได้ กฎไม่ได้ใช้ anomaly score เดี่ยว แต่ใช้หลาย field ประกอบกัน แล้วมี hard-negative filter คอยดึง false positive ออก")
    add_panel(
        slide,
        7.88,
        0.31,
        10.70,
        7.13,
        "PIPELINE ที่ใช้กับ offline_02",
        [
            ("01  NORMAL PROFILE", "เก็บชุด window และ packet shape ที่พบใน Normal จริง"),
            ("02  PROTOCOL RULES", "อ่าน SYN, CONNECT, SUBSCRIBE และ PUBLISH เป็นกลุ่ม"),
            ("03  CONTEXT GATE", "ดู tcp.stream และ family ไม่ตัดสินจาก packet เดี่ยว"),
            ("04  HARD NEGATIVE", "ตัด PINGRESP ที่ดูแปลกแต่เป็น Normal ออกจากผลบวก"),
        ],
        footer="RULES ที่ตรวจสอบได้  >  anomaly score เดี่ยว",
    )

    # Slide 7: AI and human validation
    slide = prs.slides[6]
    style_text(slide.shapes[2], "AI + HUMAN\nVALIDATION", 31, bold=True, valign=MSO_ANCHOR.MIDDLE)
    style_text(slide.shapes[3], "AI ช่วย\n• อธิบาย feature\n• เสนอ rule และ code\n• จัดอันดับ residual rows\n• ทดลอง ExtraTrees / RandomForest /\n  HistGradientBoosting / IsolationForest\n\nคนตัดสิน\n• อ่าน protocol\n• เช็ก stream context\n• หักล้าง false positive", 13.2, bold=True, line_spacing=1.02)
    narration(slide, 4, "AI ช่วยให้ผมตั้งสมมติฐานและทดลองได้เร็วขึ้น แต่ผมไม่ใช้ผลจาก AI ครั้งเดียวแล้วเชื่อทันทีครับ ทุกแนวคิดต้องกลับมาตรวจด้วย protocol meaning, stream context และ hard negatives โดยเฉพาะ ML models ที่ใช้ pseudo-label เพราะมีความเสี่ยงเรื่อง circularity")
    style_text(slide.shapes[5], "AI + EVIDENCE + HUMAN JUDGMENT", 11.5, bold=True, color=MUTED, valign=MSO_ANCHOR.MIDDLE)

    # Slide 8: selected result
    slide = prs.slides[7]
    style_text(slide.shapes[2], "THE RESULT\nWE EXPLAIN", 31, bold=True, valign=MSO_ANCHOR.MIDDLE)
    style_text(slide.shapes[3], "submission_offline_02_\nhedge_payload132.csv\n\nROWS                 10,000\nPOSITIVE LABELS        2,811\nPRIVATE F1          0.97198\nPUBLIC F1           0.96267\n\nBASELINE + STRUCTURAL Id=9816\n+ PAYLOAD132 HEDGE Id=1145", 13.7, bold=True, line_spacing=1.03)
    narration(slide, 4, "ไฟล์ที่เลือกสำหรับอธิบาย\n\nไฟล์ที่ผมนำมาอธิบายในวันนี้คือ offline_02 ครับ ไฟล์นี้ต่อยอดจาก baseline ด้วย structural row ที่ช่วยเติมรูปแบบ SYN stream และ payload hedge ที่ Id 1145 โดยตั้งใจเพิ่มเฉพาะกลุ่มที่มีเหตุผล ไม่ได้กวาดทุกแถวที่ดูผิดปกติ\n\nผลที่รายงานคือ Private F1 0.97198 และ Public F1 0.96267")
    # These two template labels sit outside the left canvas on this slide.
    style_text(slide.shapes[5], "", 11.5, bold=True, color=MUTED, valign=MSO_ANCHOR.MIDDLE)
    style_text(slide.shapes[6], "", 11.5, bold=True, color=MUTED, valign=MSO_ANCHOR.MIDDLE)
    add_panel(
        slide,
        7.88,
        0.15,
        11.01,
        7.35,
        "SELECTED FILE  |  offline_02",
        [
            ("FILE", "submission_offline_02_hedge_payload132.csv"),
            ("SCORE", "PRIVATE F1   0.97198\nPUBLIC F1    0.96267"),
            ("DATA", "10,000 rows  |  2,811 positive labels"),
            ("CHANGE", "baseline + structural Id=9816\n+ payload hedge Id=1145"),
        ],
        footer="ตัวเลขเป็นคะแนนที่ผู้ใช้รายงานจาก Kaggle",
    )

    # Slide 9: lessons
    slide = prs.slides[8]
    style_text(slide.shapes[2], "WHAT WE\nLEARNED", 31, bold=True, valign=MSO_ANCHOR.MIDDLE)
    style_text(slide.shapes[3], "1  Rare ไม่ได้แปลว่า Attack\n2  Context สำคัญกว่า packet เดี่ยว\n3  ทำนาย Attack มากไป F1 อาจลด\n4  ML score ไม่ใช่หลักฐาน\n5  กฎง่าย แต่ตรวจสอบได้", 15.2, bold=True, line_spacing=1.08)
    narration(slide, 4, "สิ่งที่ได้เรียนรู้\n\nบทเรียนที่สำคัญที่สุดคือ anomaly ไม่ได้แปลว่า attack เสมอไปครับ Normal รูปแบบใหม่ก็แตกต่างจาก Train ได้เหมือนกัน ดังนั้นระบบที่ดีไม่ใช่ระบบที่ทำนาย Attack เยอะที่สุด แต่เป็นระบบที่รู้ว่าเมื่อไหร่ควรเชื่อ และเมื่อไหร่ควรหยุด")
    style_text(slide.shapes[5], "", 15, bold=True, valign=MSO_ANCHOR.MIDDLE)

    # Slide 10: conclusion
    slide = prs.slides[9]
    style_text(slide.shapes[2], "FINAL\nTAKEAWAY", 31, bold=True, valign=MSO_ANCHOR.MIDDLE)
    style_text(slide.shapes[3], "Normal-only data\n        ↓\nLearn protocol behavior\n        ↓\nDetect meaningful deviations\n        ↓\nValidate with evidence\n\nไฟล์ที่นำมาอธิบาย\noffline_02\nPrivate 0.97198\nPublic  0.96267", 14.3, bold=True, line_spacing=1.02)
    narration(slide, 4, "บทสรุป\n\nสรุปคือผมเปลี่ยนข้อจำกัดของข้อมูล Normal-only ให้เป็นวิธีคิดแบบ protocol-first ครับ เริ่มจากเรียนรู้พฤติกรรมปกติ ใช้ TCP และ MQTT เป็นหลักฐาน เพิ่ม context ระดับ stream และใช้ hard negatives ป้องกัน false positive นี่คือเหตุผลที่ไฟล์ offline_02 อธิบายได้ทั้งวิธีทำและเหตุผลของผลลัพธ์")
    style_text(slide.shapes[5], "", 15, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_panel(
        slide,
        7.86,
        0.16,
        11.02,
        7.35,
        "PROTOCOL-FIRST SUMMARY",
        [
            ("01  NORMAL-ONLY", "เรียนรู้ขอบเขตจาก Train ที่มี Normal เท่านั้น"),
            ("02  TCP + MQTT", "ใช้ feature ที่มีความหมายของ protocol ไม่ใช่ค่าที่แปลกอย่างเดียว"),
            ("03  CONTEXT + NEGATIVE", "ดู stream และหักล้าง false positive ก่อนตัดสิน"),
            ("04  RESULT", "offline_02  |  Private F1 0.97198  |  Public F1 0.96267"),
        ],
        footer="เรียนรู้  ->  ตรวจหลักฐาน  ->  เลือกเฉพาะ deviation ที่มีเหตุผล",
    )

    # Slide 11: close
    slide = prs.slides[10]
    style_text(slide.shapes[3], "THANK YOU", 34, bold=True, valign=MSO_ANCHOR.MIDDLE)
    style_text(slide.shapes[4], "FOR LISTENING\nQUESTIONS?", 18, bold=True, color=MUTED, valign=MSO_ANCHOR.MIDDLE)

    prs.core_properties.title = "IoT Attack Detection through TCP and MQTT"
    prs.core_properties.subject = "Technical presentation for submission_offline_02_hedge_payload132.csv"
    prs.core_properties.author = "Wachirawit 610686"
    prs.core_properties.comments = "Built from the Black and Gray Gradient Professional Presentation template. Scores are user-reported Kaggle results."
    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    build()
